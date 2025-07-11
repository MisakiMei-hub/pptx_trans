from pptx import Presentation
from deep_translator import GoogleTranslator

# 获取输入文件名（自动加 .pptx 后缀）
input_name = input("请输入要翻译的 PPTX 文件名（不含扩展名）：")
pptx_file = input_name + ".pptx"
output_file = input_name + "_trans.pptx"

prs = Presentation(pptx_file)
translator = GoogleTranslator(source='zh-TW', target='zh-CN')

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    original_text = run.text
                    if isinstance(original_text, str) and original_text.strip():
                        try:
                            translated = translator.translate(original_text)
                            run.text = str(translated)
                        except Exception as e:
                            print(f"翻译出错：{original_text} -> {e}")


prs.save(output_file)
print("✅ 翻译完成，已保存为{output_file}")

